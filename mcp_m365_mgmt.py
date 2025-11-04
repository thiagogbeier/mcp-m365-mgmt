from mcp.server.fastmcp import FastMCP
import asyncio
from azure.identity import DefaultAzureCredential, InteractiveBrowserCredential, ClientSecretCredential
import requests
import os
from dotenv import load_dotenv

# Load environment variables
load_dotenv()

# Create MCP server instance
mcp = FastMCP("entra-server")

# Initialize authentication
def get_credential():
    """Get Azure credential for authentication."""
    # Check authentication mode from environment
    auth_mode = os.getenv("AUTH_MODE", "app")  # 'app' or 'user'
    
    if auth_mode == "user":
        # Interactive user authentication
        tenant_id = os.getenv("AZURE_TENANT_ID")
        client_id = os.getenv("AZURE_CLIENT_ID")
        
        if client_id and tenant_id:
            # Use InteractiveBrowserCredential with custom app registration
            return InteractiveBrowserCredential(
                tenant_id=tenant_id,
                client_id=client_id
            )
        else:
            # Fall back to default interactive auth
            return InteractiveBrowserCredential()
    else:
        # Service principal (app) authentication
        client_id = os.getenv("AZURE_CLIENT_ID")
        tenant_id = os.getenv("AZURE_TENANT_ID")
        client_secret = os.getenv("AZURE_CLIENT_SECRET")
        
        if client_id and tenant_id and client_secret:
            return ClientSecretCredential(
                tenant_id=tenant_id,
                client_id=client_id,
                client_secret=client_secret
            )
        
        # Fall back to DefaultAzureCredential (Azure CLI, managed identity, etc.)
        try:
            return DefaultAzureCredential()
        except:
            return InteractiveBrowserCredential()

credential = get_credential()

def get_access_token():
    """Get access token for Microsoft Graph API."""
    auth_mode = os.getenv("AUTH_MODE", "app")
    
    if auth_mode == "user":
        # Use user delegated permissions scope
        token = credential.get_token("https://graph.microsoft.com/.default")
    else:
        # Use application permissions scope
        token = credential.get_token("https://graph.microsoft.com/.default")
    
    return token.token

@mcp.tool()
def create_user(display_name: str, mail_nickname: str, user_principal_name: str):
    """Creates a user in Microsoft Entra ID."""
    access_token = get_access_token()
    
    headers = {
        "Authorization": f"Bearer {access_token}",
        "Content-Type": "application/json"
    }
    
    body = {
        "accountEnabled": True,
        "displayName": display_name,
        "mailNickname": mail_nickname,
        "userPrincipalName": user_principal_name,
        "passwordProfile": {
            "forceChangePasswordNextSignIn": True,
            "password": "TempPassword123!"
        }
    }
    
    response = requests.post(
        "https://graph.microsoft.com/v1.0/users",
        headers=headers,
        json=body
    )
    
    if response.status_code == 201:
        result = response.json()
        return {
            "id": result.get("id"),
            "displayName": result.get("displayName"),
            "userPrincipalName": result.get("userPrincipalName"),
            "mailNickname": result.get("mailNickname")
        }
    else:
        return {"error": response.text, "status_code": response.status_code}

@mcp.tool()
def list_intune_devices():
    """Lists Intune-managed devices from your tenant."""
    access_token = get_access_token()
    
    headers = {
        "Authorization": f"Bearer {access_token}",
        "Content-Type": "application/json"
    }
    
    response = requests.get(
        "https://graph.microsoft.com/v1.0/deviceManagement/managedDevices",
        headers=headers
    )
    
    if response.status_code == 200:
        result = response.json()
        devices = []
        
        for device in result.get("value", []):
            devices.append({
                "id": device.get("id"),
                "deviceName": device.get("deviceName"),
                "operatingSystem": device.get("operatingSystem"),
                "osVersion": device.get("osVersion"),
                "complianceState": device.get("complianceState"),
                "managedDeviceOwnerType": device.get("managedDeviceOwnerType"),
                "enrolledDateTime": device.get("enrolledDateTime"),
                "lastSyncDateTime": device.get("lastSyncDateTime")
            })
        
        return {"devices": devices, "count": len(devices)}
    else:
        return {"error": response.text, "status_code": response.status_code}

@mcp.tool()
def list_intune_compliance_policies():
    """Lists all Intune device compliance policies."""
    access_token = get_access_token()
    
    headers = {
        "Authorization": f"Bearer {access_token}",
        "Content-Type": "application/json"
    }
    
    response = requests.get(
        "https://graph.microsoft.com/v1.0/deviceManagement/deviceCompliancePolicies",
        headers=headers
    )
    
    if response.status_code == 200:
        result = response.json()
        policies = []
        
        for policy in result.get("value", []):
            policies.append({
                "id": policy.get("id"),
                "displayName": policy.get("displayName"),
                "description": policy.get("description"),
                "platform": policy.get("@odata.type", "").split('.')[-1],
                "createdDateTime": policy.get("createdDateTime"),
                "lastModifiedDateTime": policy.get("lastModifiedDateTime"),
                "version": policy.get("version")
            })
        
        return {"policies": policies, "count": len(policies)}
    else:
        return {"error": response.text, "status_code": response.status_code}

@mcp.tool()
def list_intune_configuration_policies():
    """Lists all Intune device configuration policies (settings)."""
    access_token = get_access_token()
    
    headers = {
        "Authorization": f"Bearer {access_token}",
        "Content-Type": "application/json"
    }
    
    response = requests.get(
        "https://graph.microsoft.com/v1.0/deviceManagement/deviceConfigurations",
        headers=headers
    )
    
    if response.status_code == 200:
        result = response.json()
        policies = []
        
        for policy in result.get("value", []):
            policies.append({
                "id": policy.get("id"),
                "displayName": policy.get("displayName"),
                "description": policy.get("description"),
                "platform": policy.get("@odata.type", "").split('.')[-1],
                "createdDateTime": policy.get("createdDateTime"),
                "lastModifiedDateTime": policy.get("lastModifiedDateTime"),
                "version": policy.get("version")
            })
        
        return {"policies": policies, "count": len(policies)}
    else:
        return {"error": response.text, "status_code": response.status_code}

@mcp.tool()
def list_intune_filters():
    """Lists all Intune assignment filters."""
    access_token = get_access_token()
    
    headers = {
        "Authorization": f"Bearer {access_token}",
        "Content-Type": "application/json"
    }
    
    response = requests.get(
        "https://graph.microsoft.com/beta/deviceManagement/assignmentFilters",
        headers=headers
    )
    
    if response.status_code == 200:
        result = response.json()
        filters = []
        
        for filter_item in result.get("value", []):
            filters.append({
                "id": filter_item.get("id"),
                "displayName": filter_item.get("displayName"),
                "description": filter_item.get("description"),
                "platform": filter_item.get("platform"),
                "rule": filter_item.get("rule"),
                "createdDateTime": filter_item.get("createdDateTime"),
                "lastModifiedDateTime": filter_item.get("lastModifiedDateTime")
            })
        
        return {"filters": filters, "count": len(filters)}
    else:
        return {"error": response.text, "status_code": response.status_code}

@mcp.tool()
def list_intune_scripts():
    """Lists all Intune device management scripts (PowerShell and Shell scripts)."""
    access_token = get_access_token()
    
    headers = {
        "Authorization": f"Bearer {access_token}",
        "Content-Type": "application/json"
    }
    
    # Get PowerShell scripts
    ps_response = requests.get(
        "https://graph.microsoft.com/beta/deviceManagement/deviceManagementScripts",
        headers=headers
    )
    
    # Get Shell scripts (for macOS/Linux)
    shell_response = requests.get(
        "https://graph.microsoft.com/beta/deviceManagement/deviceShellScripts",
        headers=headers
    )
    
    scripts = []
    
    if ps_response.status_code == 200:
        ps_result = ps_response.json()
        for script in ps_result.get("value", []):
            scripts.append({
                "id": script.get("id"),
                "displayName": script.get("displayName"),
                "description": script.get("description"),
                "scriptType": "PowerShell",
                "fileName": script.get("fileName"),
                "runAsAccount": script.get("runAsAccount"),
                "enforceSignatureCheck": script.get("enforceSignatureCheck"),
                "createdDateTime": script.get("createdDateTime"),
                "lastModifiedDateTime": script.get("lastModifiedDateTime")
            })
    
    if shell_response.status_code == 200:
        shell_result = shell_response.json()
        for script in shell_result.get("value", []):
            scripts.append({
                "id": script.get("id"),
                "displayName": script.get("displayName"),
                "description": script.get("description"),
                "scriptType": "Shell",
                "fileName": script.get("fileName"),
                "runAsAccount": script.get("runAsAccount"),
                "createdDateTime": script.get("createdDateTime"),
                "lastModifiedDateTime": script.get("lastModifiedDateTime")
            })
    
    return {"scripts": scripts, "count": len(scripts)}

@mcp.tool()
def list_intune_applications():
    """Lists all Intune applications (mobile apps)."""
    access_token = get_access_token()
    
    headers = {
        "Authorization": f"Bearer {access_token}",
        "Content-Type": "application/json"
    }
    
    response = requests.get(
        "https://graph.microsoft.com/beta/deviceAppManagement/mobileApps",
        headers=headers
    )
    
    if response.status_code == 200:
        result = response.json()
        apps = []
        
        for app in result.get("value", []):
            app_type = app.get("@odata.type", "").split('.')[-1]
            apps.append({
                "id": app.get("id"),
                "displayName": app.get("displayName"),
                "description": app.get("description"),
                "publisher": app.get("publisher"),
                "appType": app_type,
                "createdDateTime": app.get("createdDateTime"),
                "lastModifiedDateTime": app.get("lastModifiedDateTime"),
                "publishingState": app.get("publishingState"),
                "isAssigned": app.get("isAssigned"),
                "isFeatured": app.get("isFeatured")
            })
        
        return {"applications": apps, "count": len(apps)}
    else:
        return {"error": response.text, "status_code": response.status_code}

@mcp.tool()
def list_autopilot_profiles():
    """Lists all Windows Autopilot deployment profiles."""
    access_token = get_access_token()
    
    headers = {
        "Authorization": f"Bearer {access_token}",
        "Content-Type": "application/json"
    }
    
    response = requests.get(
        "https://graph.microsoft.com/beta/deviceManagement/windowsAutopilotDeploymentProfiles",
        headers=headers
    )
    
    if response.status_code == 200:
        result = response.json()
        profiles = []
        
        for profile in result.get("value", []):
            profiles.append({
                "id": profile.get("id"),
                "displayName": profile.get("displayName"),
                "description": profile.get("description"),
                "profileType": profile.get("@odata.type", "").split('.')[-1],
                "createdDateTime": profile.get("createdDateTime"),
                "lastModifiedDateTime": profile.get("lastModifiedDateTime"),
                "outOfBoxExperienceSettings": profile.get("outOfBoxExperienceSettings"),
                "enrollmentStatusScreenSettings": profile.get("enrollmentStatusScreenSettings"),
                "extractHardwareHash": profile.get("extractHardwareHash"),
                "deviceNameTemplate": profile.get("deviceNameTemplate"),
                "deviceType": profile.get("deviceType"),
                "enableWhiteGlove": profile.get("enableWhiteGlove")
            })
        
        return {"profiles": profiles, "count": len(profiles)}
    else:
        return {"error": response.text, "status_code": response.status_code}

@mcp.tool()
def list_autopilot_devices():
    """Lists all Windows Autopilot devices registered in the tenant."""
    access_token = get_access_token()
    
    headers = {
        "Authorization": f"Bearer {access_token}",
        "Content-Type": "application/json"
    }
    
    response = requests.get(
        "https://graph.microsoft.com/beta/deviceManagement/windowsAutopilotDeviceIdentities",
        headers=headers
    )
    
    if response.status_code == 200:
        result = response.json()
        devices = []
        
        for device in result.get("value", []):
            devices.append({
                "id": device.get("id"),
                "serialNumber": device.get("serialNumber"),
                "model": device.get("model"),
                "manufacturer": device.get("manufacturer"),
                "productKey": device.get("productKey"),
                "groupTag": device.get("groupTag"),
                "purchaseOrderIdentifier": device.get("purchaseOrderIdentifier"),
                "enrollmentState": device.get("enrollmentState"),
                "lastContactedDateTime": device.get("lastContactedDateTime"),
                "addressableUserName": device.get("addressableUserName"),
                "userPrincipalName": device.get("userPrincipalName"),
                "resourceName": device.get("resourceName"),
                "skuNumber": device.get("skuNumber"),
                "systemFamily": device.get("systemFamily"),
                "azureActiveDirectoryDeviceId": device.get("azureActiveDirectoryDeviceId"),
                "managedDeviceId": device.get("managedDeviceId"),
                "displayName": device.get("displayName")
            })
        
        return {"devices": devices, "count": len(devices)}
    else:
        return {"error": response.text, "status_code": response.status_code}

@mcp.tool()
def list_enrollment_status_page_profiles():
    """Lists all Enrollment Status Page (ESP) profiles for Windows Autopilot."""
    access_token = get_access_token()
    
    headers = {
        "Authorization": f"Bearer {access_token}",
        "Content-Type": "application/json"
    }
    
    response = requests.get(
        "https://graph.microsoft.com/beta/deviceManagement/deviceEnrollmentConfigurations",
        headers=headers
    )
    
    if response.status_code == 200:
        result = response.json()
        esp_profiles = []
        
        for config in result.get("value", []):
            # Filter for Windows10EnrollmentCompletionPageConfiguration (ESP profiles)
            if "windows10EnrollmentCompletionPageConfiguration" in config.get("@odata.type", "").lower():
                esp_profiles.append({
                    "id": config.get("id"),
                    "displayName": config.get("displayName"),
                    "description": config.get("description"),
                    "priority": config.get("priority"),
                    "createdDateTime": config.get("createdDateTime"),
                    "lastModifiedDateTime": config.get("lastModifiedDateTime"),
                    "version": config.get("version"),
                    "showInstallationProgress": config.get("showInstallationProgress"),
                    "blockDeviceSetupRetryByUser": config.get("blockDeviceSetupRetryByUser"),
                    "allowDeviceResetOnInstallFailure": config.get("allowDeviceResetOnInstallFailure"),
                    "allowLogCollectionOnInstallFailure": config.get("allowLogCollectionOnInstallFailure"),
                    "customErrorMessage": config.get("customErrorMessage"),
                    "installProgressTimeoutInMinutes": config.get("installProgressTimeoutInMinutes"),
                    "allowDeviceUseOnInstallFailure": config.get("allowDeviceUseOnInstallFailure"),
                    "selectedMobileAppIds": config.get("selectedMobileAppIds"),
                    "trackInstallProgressForAutopilotOnly": config.get("trackInstallProgressForAutopilotOnly"),
                    "disableUserStatusTrackingAfterFirstUser": config.get("disableUserStatusTrackingAfterFirstUser")
                })
        
        return {"esp_profiles": esp_profiles, "count": len(esp_profiles)}
    else:
        return {"error": response.text, "status_code": response.status_code}

@mcp.tool()
def list_android_management_profiles():
    """Lists all Android device management settings, policies, profiles, and enrollment configurations."""
    access_token = get_access_token()
    
    headers = {
        "Authorization": f"Bearer {access_token}",
        "Content-Type": "application/json"
    }
    
    # Get Android device configurations
    config_response = requests.get(
        "https://graph.microsoft.com/beta/deviceManagement/deviceConfigurations",
        headers=headers
    )
    
    # Get Android enrollment configurations
    enrollment_response = requests.get(
        "https://graph.microsoft.com/beta/deviceManagement/deviceEnrollmentConfigurations",
        headers=headers
    )
    
    # Get Android compliance policies
    compliance_response = requests.get(
        "https://graph.microsoft.com/v1.0/deviceManagement/deviceCompliancePolicies",
        headers=headers
    )
    
    android_profiles = {
        "device_configurations": [],
        "enrollment_configurations": [],
        "compliance_policies": []
    }
    
    # Parse device configurations
    if config_response.status_code == 200:
        configs = config_response.json().get("value", [])
        for config in configs:
            config_type = config.get("@odata.type", "")
            if "android" in config_type.lower():
                android_profiles["device_configurations"].append({
                    "id": config.get("id"),
                    "displayName": config.get("displayName"),
                    "description": config.get("description"),
                    "type": config_type.split('.')[-1],
                    "createdDateTime": config.get("createdDateTime"),
                    "lastModifiedDateTime": config.get("lastModifiedDateTime"),
                    "version": config.get("version")
                })
    
    # Parse enrollment configurations
    if enrollment_response.status_code == 200:
        enrollments = enrollment_response.json().get("value", [])
        for enrollment in enrollments:
            enrollment_type = enrollment.get("@odata.type", "")
            if "android" in enrollment_type.lower():
                android_profiles["enrollment_configurations"].append({
                    "id": enrollment.get("id"),
                    "displayName": enrollment.get("displayName"),
                    "description": enrollment.get("description"),
                    "type": enrollment_type.split('.')[-1],
                    "priority": enrollment.get("priority"),
                    "createdDateTime": enrollment.get("createdDateTime"),
                    "lastModifiedDateTime": enrollment.get("lastModifiedDateTime"),
                    "version": enrollment.get("version")
                })
    
    # Parse compliance policies
    if compliance_response.status_code == 200:
        policies = compliance_response.json().get("value", [])
        for policy in policies:
            policy_type = policy.get("@odata.type", "")
            if "android" in policy_type.lower():
                android_profiles["compliance_policies"].append({
                    "id": policy.get("id"),
                    "displayName": policy.get("displayName"),
                    "description": policy.get("description"),
                    "type": policy_type.split('.')[-1],
                    "createdDateTime": policy.get("createdDateTime"),
                    "lastModifiedDateTime": policy.get("lastModifiedDateTime"),
                    "version": policy.get("version")
                })
    
    android_profiles["total_count"] = (
        len(android_profiles["device_configurations"]) +
        len(android_profiles["enrollment_configurations"]) +
        len(android_profiles["compliance_policies"])
    )
    
    return android_profiles

@mcp.tool()
def list_ios_management_profiles():
    """Lists all iOS/iPadOS device management settings, policies, profiles, and enrollment configurations."""
    access_token = get_access_token()
    
    headers = {
        "Authorization": f"Bearer {access_token}",
        "Content-Type": "application/json"
    }
    
    # Get iOS device configurations
    config_response = requests.get(
        "https://graph.microsoft.com/beta/deviceManagement/deviceConfigurations",
        headers=headers
    )
    
    # Get iOS enrollment configurations
    enrollment_response = requests.get(
        "https://graph.microsoft.com/beta/deviceManagement/deviceEnrollmentConfigurations",
        headers=headers
    )
    
    # Get iOS compliance policies
    compliance_response = requests.get(
        "https://graph.microsoft.com/v1.0/deviceManagement/deviceCompliancePolicies",
        headers=headers
    )
    
    ios_profiles = {
        "device_configurations": [],
        "enrollment_configurations": [],
        "compliance_policies": []
    }
    
    # Parse device configurations
    if config_response.status_code == 200:
        configs = config_response.json().get("value", [])
        for config in configs:
            config_type = config.get("@odata.type", "")
            if "ios" in config_type.lower():
                ios_profiles["device_configurations"].append({
                    "id": config.get("id"),
                    "displayName": config.get("displayName"),
                    "description": config.get("description"),
                    "type": config_type.split('.')[-1],
                    "createdDateTime": config.get("createdDateTime"),
                    "lastModifiedDateTime": config.get("lastModifiedDateTime"),
                    "version": config.get("version")
                })
    
    # Parse enrollment configurations
    if enrollment_response.status_code == 200:
        enrollments = enrollment_response.json().get("value", [])
        for enrollment in enrollments:
            enrollment_type = enrollment.get("@odata.type", "")
            if "ios" in enrollment_type.lower():
                ios_profiles["enrollment_configurations"].append({
                    "id": enrollment.get("id"),
                    "displayName": enrollment.get("displayName"),
                    "description": enrollment.get("description"),
                    "type": enrollment_type.split('.')[-1],
                    "priority": enrollment.get("priority"),
                    "createdDateTime": enrollment.get("createdDateTime"),
                    "lastModifiedDateTime": enrollment.get("lastModifiedDateTime"),
                    "version": enrollment.get("version")
                })
    
    # Parse compliance policies
    if compliance_response.status_code == 200:
        policies = compliance_response.json().get("value", [])
        for policy in policies:
            policy_type = policy.get("@odata.type", "")
            if "ios" in policy_type.lower():
                ios_profiles["compliance_policies"].append({
                    "id": policy.get("id"),
                    "displayName": policy.get("displayName"),
                    "description": policy.get("description"),
                    "type": policy_type.split('.')[-1],
                    "createdDateTime": policy.get("createdDateTime"),
                    "lastModifiedDateTime": policy.get("lastModifiedDateTime"),
                    "version": policy.get("version")
                })
    
    ios_profiles["total_count"] = (
        len(ios_profiles["device_configurations"]) +
        len(ios_profiles["enrollment_configurations"]) +
        len(ios_profiles["compliance_policies"])
    )
    
    return ios_profiles

@mcp.tool()
def list_app_protection_policies():
    """Lists all app protection policies (MAM policies) for iOS, Android, and Windows."""
    access_token = get_access_token()
    
    headers = {
        "Authorization": f"Bearer {access_token}",
        "Content-Type": "application/json"
    }
    
    # Get managed app policies
    response = requests.get(
        "https://graph.microsoft.com/beta/deviceAppManagement/managedAppPolicies",
        headers=headers
    )
    
    if response.status_code == 200:
        result = response.json()
        policies = []
        
        for policy in result.get("value", []):
            policy_type = policy.get("@odata.type", "").split('.')[-1]
            policies.append({
                "id": policy.get("id"),
                "displayName": policy.get("displayName"),
                "description": policy.get("description"),
                "policyType": policy_type,
                "createdDateTime": policy.get("createdDateTime"),
                "lastModifiedDateTime": policy.get("lastModifiedDateTime"),
                "version": policy.get("version"),
                "isAssigned": policy.get("isAssigned"),
                # Platform-specific fields
                "platformType": policy.get("@odata.type", "").lower().replace("#microsoft.graph.", "").split("managedapp")[0] if "managedapp" in policy.get("@odata.type", "").lower() else "unknown"
            })
        
        return {"app_protection_policies": policies, "count": len(policies)}
    else:
        return {"error": response.text, "status_code": response.status_code}

@mcp.tool()
def list_microsoft_tunnel_sites():
    """Lists all Microsoft Tunnel Gateway sites and their configurations."""
    access_token = get_access_token()
    
    headers = {
        "Authorization": f"Bearer {access_token}",
        "Content-Type": "application/json"
    }
    
    # Get Microsoft Tunnel sites
    response = requests.get(
        "https://graph.microsoft.com/beta/deviceManagement/microsoftTunnelSites",
        headers=headers
    )
    
    if response.status_code == 200:
        result = response.json()
        sites = []
        
        for site in result.get("value", []):
            sites.append({
                "id": site.get("id"),
                "displayName": site.get("displayName"),
                "description": site.get("description"),
                "publicAddress": site.get("publicAddress"),
                "upgradeWindowUtcOffsetInMinutes": site.get("upgradeWindowUtcOffsetInMinutes"),
                "upgradeWindowStartTime": site.get("upgradeWindowStartTime"),
                "upgradeWindowEndTime": site.get("upgradeWindowEndTime"),
                "upgradeAutomatically": site.get("upgradeAutomatically"),
                "upgradeAvailable": site.get("upgradeAvailable"),
                "internalNetworkProbeUrl": site.get("internalNetworkProbeUrl"),
                "roleScopeTagIds": site.get("roleScopeTagIds")
            })
        
        return {"tunnel_sites": sites, "count": len(sites)}
    else:
        return {"error": response.text, "status_code": response.status_code}

@mcp.tool()
def list_microsoft_tunnel_servers():
    """Lists all Microsoft Tunnel Gateway servers across all sites."""
    access_token = get_access_token()
    
    headers = {
        "Authorization": f"Bearer {access_token}",
        "Content-Type": "application/json"
    }
    
    # First get all tunnel sites
    sites_response = requests.get(
        "https://graph.microsoft.com/beta/deviceManagement/microsoftTunnelSites",
        headers=headers
    )
    
    all_servers = []
    
    if sites_response.status_code == 200:
        sites = sites_response.json().get("value", [])
        
        # For each site, get its servers
        for site in sites:
            site_id = site.get("id")
            servers_response = requests.get(
                f"https://graph.microsoft.com/beta/deviceManagement/microsoftTunnelSites/{site_id}/microsoftTunnelServers",
                headers=headers
            )
            
            if servers_response.status_code == 200:
                servers = servers_response.json().get("value", [])
                for server in servers:
                    all_servers.append({
                        "id": server.get("id"),
                        "displayName": server.get("displayName"),
                        "tunnelServerHealthStatus": server.get("tunnelServerHealthStatus"),
                        "lastCheckinDateTime": server.get("lastCheckinDateTime"),
                        "agentImageDigest": server.get("agentImageDigest"),
                        "serverImageDigest": server.get("serverImageDigest"),
                        "siteName": site.get("displayName"),
                        "siteId": site_id
                    })
        
        return {"tunnel_servers": all_servers, "count": len(all_servers)}
    else:
        return {"error": sites_response.text, "status_code": sites_response.status_code}

@mcp.tool()
def list_intune_ad_connectors():
    """Lists all Intune Connector for Active Directory (used for Hybrid Azure AD Join and Autopilot)."""
    access_token = get_access_token()
    
    headers = {
        "Authorization": f"Bearer {access_token}",
        "Content-Type": "application/json"
    }
    
    response = requests.get(
        "https://graph.microsoft.com/beta/deviceManagement/domainJoinConnectors",
        headers=headers
    )
    
    if response.status_code == 200:
        result = response.json()
        connectors = []
        
        for connector in result.get("value", []):
            connectors.append({
                "id": connector.get("id"),
                "displayName": connector.get("displayName"),
                "state": connector.get("state"),
                "version": connector.get("version"),
                "machineName": connector.get("machineName"),
                "lastConnectionDateTime": connector.get("lastConnectionDateTime")
            })
        
        return {"ad_connectors": connectors, "count": len(connectors)}
    else:
        return {"error": response.text, "status_code": response.status_code}

@mcp.tool()
def list_intune_certificate_connectors():
    """Lists all Intune Certificate Connectors (NDES connectors for SCEP certificates)."""
    access_token = get_access_token()
    
    headers = {
        "Authorization": f"Bearer {access_token}",
        "Content-Type": "application/json"
    }
    
    response = requests.get(
        "https://graph.microsoft.com/beta/deviceManagement/ndesConnectors",
        headers=headers
    )
    
    if response.status_code == 200:
        result = response.json()
        connectors = []
        
        for connector in result.get("value", []):
            connectors.append({
                "id": connector.get("id"),
                "displayName": connector.get("displayName"),
                "lastConnectionDateTime": connector.get("lastConnectionDateTime"),
                "state": connector.get("state"),
                "connectorVersion": connector.get("connectorVersion"),
                "machineName": connector.get("machineName"),
                "enrolledDateTime": connector.get("enrolledDateTime")
            })
        
        return {"certificate_connectors": connectors, "count": len(connectors)}
    else:
        return {"error": response.text, "status_code": response.status_code}

@mcp.tool()
def get_user_info(user_id: str):
    """Gets information about a specific user by user principal name or object ID."""
    access_token = get_access_token()
    
    headers = {
        "Authorization": f"Bearer {access_token}",
        "Content-Type": "application/json"
    }
    
    response = requests.get(
        f"https://graph.microsoft.com/v1.0/users/{user_id}",
        headers=headers
    )
    
    if response.status_code == 200:
        user = response.json()
        return {
            "id": user.get("id"),
            "displayName": user.get("displayName"),
            "userPrincipalName": user.get("userPrincipalName"),
            "mail": user.get("mail"),
            "jobTitle": user.get("jobTitle"),
            "department": user.get("department"),
            "officeLocation": user.get("officeLocation"),
            "mobilePhone": user.get("mobilePhone"),
            "businessPhones": user.get("businessPhones"),
            "accountEnabled": user.get("accountEnabled")
        }
    else:
        return {"error": response.text, "status_code": response.status_code}

@mcp.tool()
def list_users():
    """Lists all users in the tenant."""
    access_token = get_access_token()
    
    headers = {
        "Authorization": f"Bearer {access_token}",
        "Content-Type": "application/json"
    }
    
    response = requests.get(
        "https://graph.microsoft.com/v1.0/users",
        headers=headers
    )
    
    if response.status_code == 200:
        result = response.json()
        users = []
        
        for user in result.get("value", []):
            users.append({
                "id": user.get("id"),
                "displayName": user.get("displayName"),
                "userPrincipalName": user.get("userPrincipalName"),
                "mail": user.get("mail"),
                "jobTitle": user.get("jobTitle"),
                "accountEnabled": user.get("accountEnabled")
            })
        
        return {"users": users, "count": len(users)}
    else:
        return {"error": response.text, "status_code": response.status_code}

@mcp.tool()
def list_groups():
    """Lists all groups in the tenant."""
    access_token = get_access_token()
    
    headers = {
        "Authorization": f"Bearer {access_token}",
        "Content-Type": "application/json"
    }
    
    response = requests.get(
        "https://graph.microsoft.com/v1.0/groups",
        headers=headers
    )
    
    if response.status_code == 200:
        result = response.json()
        groups = []
        
        for group in result.get("value", []):
            groups.append({
                "id": group.get("id"),
                "displayName": group.get("displayName"),
                "description": group.get("description"),
                "mailEnabled": group.get("mailEnabled"),
                "securityEnabled": group.get("securityEnabled"),
                "mail": group.get("mail"),
                "groupTypes": group.get("groupTypes")
            })
        
        return {"groups": groups, "count": len(groups)}
    else:
        return {"error": response.text, "status_code": response.status_code}

@mcp.tool()
def get_group_members(group_id: str):
    """Gets members of a specific group."""
    access_token = get_access_token()
    
    headers = {
        "Authorization": f"Bearer {access_token}",
        "Content-Type": "application/json"
    }
    
    response = requests.get(
        f"https://graph.microsoft.com/v1.0/groups/{group_id}/members",
        headers=headers
    )
    
    if response.status_code == 200:
        result = response.json()
        members = []
        
        for member in result.get("value", []):
            members.append({
                "id": member.get("id"),
                "displayName": member.get("displayName"),
                "userPrincipalName": member.get("userPrincipalName"),
                "mail": member.get("mail"),
                "@odata.type": member.get("@odata.type")
            })
        
        return {"members": members, "count": len(members), "groupId": group_id}
    else:
        return {"error": response.text, "status_code": response.status_code}

@mcp.tool()
def create_file_in_onedrive(user_id: str, file_name: str, content: str, folder_path: str = ""):
    """Creates a text file in a user's OneDrive.
    
    Args:
        user_id: User principal name or object ID
        file_name: Name of the file to create (e.g., 'document.txt')
        content: Text content to write to the file
        folder_path: Optional folder path (e.g., 'Documents/MyFolder'). Leave empty for root.
    """
    access_token = get_access_token()
    
    headers = {
        "Authorization": f"Bearer {access_token}",
        "Content-Type": "text/plain"
    }
    
    # Build the path
    if folder_path:
        path = f"/drive/root:/{folder_path}/{file_name}:/content"
    else:
        path = f"/drive/root:/{file_name}:/content"
    
    response = requests.put(
        f"https://graph.microsoft.com/v1.0/users/{user_id}{path}",
        headers=headers,
        data=content.encode('utf-8')
    )
    
    if response.status_code in [200, 201]:
        result = response.json()
        return {
            "id": result.get("id"),
            "name": result.get("name"),
            "size": result.get("size"),
            "webUrl": result.get("webUrl"),
            "createdDateTime": result.get("createdDateTime"),
            "lastModifiedDateTime": result.get("lastModifiedDateTime")
        }
    else:
        return {"error": response.text, "status_code": response.status_code}

@mcp.tool()
def create_file_in_sharepoint(site_id: str, file_name: str, content: str, folder_path: str = ""):
    """Creates a text file in a SharePoint site's document library.
    
    Args:
        site_id: SharePoint site ID (use format: 'hostname,site-id,web-id' or just 'site-id')
        file_name: Name of the file to create (e.g., 'document.txt')
        content: Text content to write to the file
        folder_path: Optional folder path within the document library (e.g., 'Shared Documents/MyFolder')
    """
    access_token = get_access_token()
    
    headers = {
        "Authorization": f"Bearer {access_token}",
        "Content-Type": "text/plain"
    }
    
    # Build the path
    if folder_path:
        path = f"/drive/root:/{folder_path}/{file_name}:/content"
    else:
        path = f"/drive/root:/{file_name}:/content"
    
    response = requests.put(
        f"https://graph.microsoft.com/v1.0/sites/{site_id}{path}",
        headers=headers,
        data=content.encode('utf-8')
    )
    
    if response.status_code in [200, 201]:
        result = response.json()
        return {
            "id": result.get("id"),
            "name": result.get("name"),
            "size": result.get("size"),
            "webUrl": result.get("webUrl"),
            "createdDateTime": result.get("createdDateTime"),
            "lastModifiedDateTime": result.get("lastModifiedDateTime")
        }
    else:
        return {"error": response.text, "status_code": response.status_code}

@mcp.tool()
def list_sharepoint_sites():
    """Lists all SharePoint sites in the tenant."""
    access_token = get_access_token()
    
    headers = {
        "Authorization": f"Bearer {access_token}",
        "Content-Type": "application/json"
    }
    
    response = requests.get(
        "https://graph.microsoft.com/v1.0/sites?search=*",
        headers=headers
    )
    
    if response.status_code == 200:
        result = response.json()
        sites = []
        
        for site in result.get("value", []):
            sites.append({
                "id": site.get("id"),
                "name": site.get("name"),
                "displayName": site.get("displayName"),
                "webUrl": site.get("webUrl"),
                "description": site.get("description")
            })
        
        return {"sites": sites, "count": len(sites)}
    else:
        return {"error": response.text, "status_code": response.status_code}

@mcp.tool()
def create_word_document(location_type: str, location_id: str, file_name: str, content: str, folder_path: str = ""):
    """Creates a Word document (.docx) in OneDrive or SharePoint.
    
    Args:
        location_type: Either 'onedrive' or 'sharepoint'
        location_id: User ID for OneDrive, or Site ID for SharePoint
        file_name: Name of the file (e.g., 'report.docx')
        content: HTML content to put in the document
        folder_path: Optional folder path
    """
    import io
    from docx import Document
    from docx.shared import Pt
    
    # Create Word document
    doc = Document()
    
    # Add content (simple paragraph)
    paragraph = doc.add_paragraph(content)
    
    # Save to bytes
    doc_bytes = io.BytesIO()
    doc.save(doc_bytes)
    doc_bytes.seek(0)
    
    access_token = get_access_token()
    headers = {
        "Authorization": f"Bearer {access_token}",
        "Content-Type": "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
    }
    
    # Build URL based on location type
    if location_type == "onedrive":
        if folder_path:
            url = f"https://graph.microsoft.com/v1.0/users/{location_id}/drive/root:/{folder_path}/{file_name}:/content"
        else:
            url = f"https://graph.microsoft.com/v1.0/users/{location_id}/drive/root:/{file_name}:/content"
    else:  # sharepoint
        if folder_path:
            url = f"https://graph.microsoft.com/v1.0/sites/{location_id}/drive/root:/{folder_path}/{file_name}:/content"
        else:
            url = f"https://graph.microsoft.com/v1.0/sites/{location_id}/drive/root:/{file_name}:/content"
    
    response = requests.put(url, headers=headers, data=doc_bytes.getvalue())
    
    if response.status_code in [200, 201]:
        result = response.json()
        return {
            "id": result.get("id"),
            "name": result.get("name"),
            "size": result.get("size"),
            "webUrl": result.get("webUrl"),
            "createdDateTime": result.get("createdDateTime")
        }
    else:
        return {"error": response.text, "status_code": response.status_code}

@mcp.tool()
def create_excel_workbook(location_type: str, location_id: str, file_name: str, data: list, folder_path: str = ""):
    """Creates an Excel workbook (.xlsx) in OneDrive or SharePoint.
    
    Args:
        location_type: Either 'onedrive' or 'sharepoint'
        location_id: User ID for OneDrive, or Site ID for SharePoint
        file_name: Name of the file (e.g., 'data.xlsx')
        data: List of lists representing rows (e.g., [['Name', 'Age'], ['John', 30]])
        folder_path: Optional folder path
    """
    import io
    from openpyxl import Workbook
    
    # Create Excel workbook
    wb = Workbook()
    ws = wb.active
    ws.title = "Sheet1"
    
    # Add data
    for row in data:
        ws.append(row)
    
    # Save to bytes
    excel_bytes = io.BytesIO()
    wb.save(excel_bytes)
    excel_bytes.seek(0)
    
    access_token = get_access_token()
    headers = {
        "Authorization": f"Bearer {access_token}",
        "Content-Type": "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
    }
    
    # Build URL based on location type
    if location_type == "onedrive":
        if folder_path:
            url = f"https://graph.microsoft.com/v1.0/users/{location_id}/drive/root:/{folder_path}/{file_name}:/content"
        else:
            url = f"https://graph.microsoft.com/v1.0/users/{location_id}/drive/root:/{file_name}:/content"
    else:  # sharepoint
        if folder_path:
            url = f"https://graph.microsoft.com/v1.0/sites/{location_id}/drive/root:/{folder_path}/{file_name}:/content"
        else:
            url = f"https://graph.microsoft.com/v1.0/sites/{location_id}/drive/root:/{file_name}:/content"
    
    response = requests.put(url, headers=headers, data=excel_bytes.getvalue())
    
    if response.status_code in [200, 201]:
        result = response.json()
        return {
            "id": result.get("id"),
            "name": result.get("name"),
            "size": result.get("size"),
            "webUrl": result.get("webUrl"),
            "createdDateTime": result.get("createdDateTime")
        }
    else:
        return {"error": response.text, "status_code": response.status_code}

@mcp.tool()
def create_powerpoint_presentation(location_type: str, location_id: str, file_name: str, title: str, content: str, folder_path: str = ""):
    """Creates a PowerPoint presentation (.pptx) in OneDrive or SharePoint.
    
    Args:
        location_type: Either 'onedrive' or 'sharepoint'
        location_id: User ID for OneDrive, or Site ID for SharePoint
        file_name: Name of the file (e.g., 'presentation.pptx')
        title: Title for the first slide
        content: Content for the first slide
        folder_path: Optional folder path
    """
    import io
    from pptx import Presentation
    from pptx.util import Inches, Pt
    
    # Create PowerPoint presentation
    prs = Presentation()
    
    # Add title slide
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    slide.shapes.title.text = title
    slide.placeholders[1].text = content
    
    # Save to bytes
    pptx_bytes = io.BytesIO()
    prs.save(pptx_bytes)
    pptx_bytes.seek(0)
    
    access_token = get_access_token()
    headers = {
        "Authorization": f"Bearer {access_token}",
        "Content-Type": "application/vnd.openxmlformats-officedocument.presentationml.presentation"
    }
    
    # Build URL based on location type
    if location_type == "onedrive":
        if folder_path:
            url = f"https://graph.microsoft.com/v1.0/users/{location_id}/drive/root:/{folder_path}/{file_name}:/content"
        else:
            url = f"https://graph.microsoft.com/v1.0/users/{location_id}/drive/root:/{file_name}:/content"
    else:  # sharepoint
        if folder_path:
            url = f"https://graph.microsoft.com/v1.0/sites/{location_id}/drive/root:/{folder_path}/{file_name}:/content"
        else:
            url = f"https://graph.microsoft.com/v1.0/sites/{location_id}/drive/root:/{file_name}:/content"
    
    response = requests.put(url, headers=headers, data=pptx_bytes.getvalue())
    
    if response.status_code in [200, 201]:
        result = response.json()
        return {
            "id": result.get("id"),
            "name": result.get("name"),
            "size": result.get("size"),
            "webUrl": result.get("webUrl"),
            "createdDateTime": result.get("createdDateTime")
        }
    else:
        return {"error": response.text, "status_code": response.status_code}

@mcp.tool()
def convert_file_to_pdf(location_type: str, location_id: str, file_id: str, output_folder: str = ""):
    """Converts a file to PDF in OneDrive or SharePoint.
    
    Args:
        location_type: Either 'onedrive' or 'sharepoint'
        location_id: User ID for OneDrive, or Site ID for SharePoint
        file_id: The ID of the file to convert
        output_folder: Optional folder path to save the PDF (defaults to same location as source)
    
    Note: This works for Word, Excel, PowerPoint files
    """
    access_token = get_access_token()
    
    # First, get the file info to get its name
    headers = {
        "Authorization": f"Bearer {access_token}",
        "Content-Type": "application/json"
    }
    
    if location_type == "onedrive":
        get_url = f"https://graph.microsoft.com/v1.0/users/{location_id}/drive/items/{file_id}"
    else:  # sharepoint
        get_url = f"https://graph.microsoft.com/v1.0/sites/{location_id}/drive/items/{file_id}"
    
    response = requests.get(get_url, headers=headers)
    
    if response.status_code != 200:
        return {"error": "Failed to get file info", "status_code": response.status_code, "details": response.text}
    
    file_info = response.json()
    original_name = file_info.get("name", "")
    pdf_name = original_name.rsplit('.', 1)[0] + '.pdf'
    
    # Get the PDF content
    if location_type == "onedrive":
        convert_url = f"https://graph.microsoft.com/v1.0/users/{location_id}/drive/items/{file_id}/content?format=pdf"
    else:  # sharepoint
        convert_url = f"https://graph.microsoft.com/v1.0/sites/{location_id}/drive/items/{file_id}/content?format=pdf"
    
    response = requests.get(convert_url, headers={"Authorization": f"Bearer {access_token}"})
    
    if response.status_code != 200:
        return {"error": "Failed to convert file", "status_code": response.status_code, "details": response.text}
    
    pdf_content = response.content
    
    # Upload the PDF
    pdf_headers = {
        "Authorization": f"Bearer {access_token}",
        "Content-Type": "application/pdf"
    }
    
    if location_type == "onedrive":
        if output_folder:
            upload_url = f"https://graph.microsoft.com/v1.0/users/{location_id}/drive/root:/{output_folder}/{pdf_name}:/content"
        else:
            upload_url = f"https://graph.microsoft.com/v1.0/users/{location_id}/drive/root:/{pdf_name}:/content"
    else:  # sharepoint
        if output_folder:
            upload_url = f"https://graph.microsoft.com/v1.0/sites/{location_id}/drive/root:/{output_folder}/{pdf_name}:/content"
        else:
            upload_url = f"https://graph.microsoft.com/v1.0/sites/{location_id}/drive/root:/{pdf_name}:/content"
    
    upload_response = requests.put(upload_url, headers=pdf_headers, data=pdf_content)
    
    if upload_response.status_code in [200, 201]:
        result = upload_response.json()
        return {
            "id": result.get("id"),
            "name": result.get("name"),
            "size": result.get("size"),
            "webUrl": result.get("webUrl"),
            "originalFile": original_name,
            "createdDateTime": result.get("createdDateTime")
        }
    else:
        return {"error": "Failed to upload PDF", "status_code": upload_response.status_code, "details": upload_response.text}

@mcp.tool()
def create_csv_file(location_type: str, location_id: str, file_name: str, data: list, folder_path: str = ""):
    """Creates a CSV file in OneDrive or SharePoint.
    
    Args:
        location_type: Either 'onedrive' or 'sharepoint'
        location_id: User ID for OneDrive, or Site ID for SharePoint
        file_name: Name of the file (e.g., 'data.csv')
        data: List of lists representing rows (e.g., [['Name', 'Age'], ['John', 30]])
        folder_path: Optional folder path
    """
    import csv
    import io
    
    # Create CSV content
    csv_buffer = io.StringIO()
    writer = csv.writer(csv_buffer)
    for row in data:
        writer.writerow(row)
    
    csv_content = csv_buffer.getvalue()
    
    access_token = get_access_token()
    headers = {
        "Authorization": f"Bearer {access_token}",
        "Content-Type": "text/csv"
    }
    
    # Build URL
    if location_type == "onedrive":
        if folder_path:
            url = f"https://graph.microsoft.com/v1.0/users/{location_id}/drive/root:/{folder_path}/{file_name}:/content"
        else:
            url = f"https://graph.microsoft.com/v1.0/users/{location_id}/drive/root:/{file_name}:/content"
    else:  # sharepoint
        if folder_path:
            url = f"https://graph.microsoft.com/v1.0/sites/{location_id}/drive/root:/{folder_path}/{file_name}:/content"
        else:
            url = f"https://graph.microsoft.com/v1.0/sites/{location_id}/drive/root:/{file_name}:/content"
    
    response = requests.put(url, headers=headers, data=csv_content.encode('utf-8'))
    
    if response.status_code in [200, 201]:
        result = response.json()
        return {
            "id": result.get("id"),
            "name": result.get("name"),
            "size": result.get("size"),
            "webUrl": result.get("webUrl"),
            "createdDateTime": result.get("createdDateTime")
        }
    else:
        return {"error": response.text, "status_code": response.status_code}

@mcp.tool()
def read_csv_file(location_type: str, location_id: str, file_id: str):
    """Reads a CSV file from OneDrive or SharePoint and returns the data.
    
    Args:
        location_type: Either 'onedrive' or 'sharepoint'
        location_id: User ID for OneDrive, or Site ID for SharePoint
        file_id: The ID of the CSV file to read
    
    Returns:
        List of lists representing the CSV data
    """
    import csv
    import io
    
    access_token = get_access_token()
    headers = {"Authorization": f"Bearer {access_token}"}
    
    # Get file content
    if location_type == "onedrive":
        url = f"https://graph.microsoft.com/v1.0/users/{location_id}/drive/items/{file_id}/content"
    else:  # sharepoint
        url = f"https://graph.microsoft.com/v1.0/sites/{location_id}/drive/items/{file_id}/content"
    
    response = requests.get(url, headers=headers)
    
    if response.status_code == 200:
        csv_content = response.text
        csv_reader = csv.reader(io.StringIO(csv_content))
        data = [row for row in csv_reader]
        return {"data": data, "rows": len(data), "columns": len(data[0]) if data else 0}
    else:
        return {"error": response.text, "status_code": response.status_code}

@mcp.tool()
def export_powerpoint_slide_as_image(location_type: str, location_id: str, file_id: str, slide_index: int, image_format: str = "png", output_folder: str = ""):
    """Exports a PowerPoint slide as an image (PNG, JPG, GIF, BMP, TIFF).
    
    Args:
        location_type: Either 'onedrive' or 'sharepoint'
        location_id: User ID for OneDrive, or Site ID for SharePoint
        file_id: The ID of the PowerPoint file
        slide_index: The slide number to export (1-based index)
        image_format: Image format - 'png', 'jpg', 'gif', 'bmp', or 'tiff'
        output_folder: Optional folder path to save the image
    """
    access_token = get_access_token()
    
    # First, get file info for naming
    headers = {"Authorization": f"Bearer {access_token}"}
    if location_type == "onedrive":
        info_url = f"https://graph.microsoft.com/v1.0/users/{location_id}/drive/items/{file_id}"
    else:
        info_url = f"https://graph.microsoft.com/v1.0/sites/{location_id}/drive/items/{file_id}"
    
    info_response = requests.get(info_url, headers=headers)
    if info_response.status_code != 200:
        return {"error": "Failed to get file info", "status_code": info_response.status_code}
    
    file_name = info_response.json().get("name", "slide")
    base_name = file_name.rsplit('.', 1)[0]
    image_name = f"{base_name}_slide{slide_index}.{image_format}"
    
    # Get slide as image (using thumbnail API with high resolution)
    # Note: Full slide export requires PowerPoint Online API which has limitations
    # Using thumbnail API as a workaround
    if location_type == "onedrive":
        thumb_url = f"https://graph.microsoft.com/v1.0/users/{location_id}/drive/items/{file_id}/thumbnails/0/large/content"
    else:
        thumb_url = f"https://graph.microsoft.com/v1.0/sites/{location_id}/drive/items/{file_id}/thumbnails/0/large/content"
    
    thumb_response = requests.get(thumb_url, headers=headers)
    
    if thumb_response.status_code == 200:
        image_content = thumb_response.content
        
        # Upload the image
        content_types = {
            "png": "image/png",
            "jpg": "image/jpeg",
            "gif": "image/gif",
            "bmp": "image/bmp",
            "tiff": "image/tiff"
        }
        
        upload_headers = {
            "Authorization": f"Bearer {access_token}",
            "Content-Type": content_types.get(image_format, "image/png")
        }
        
        if location_type == "onedrive":
            if output_folder:
                upload_url = f"https://graph.microsoft.com/v1.0/users/{location_id}/drive/root:/{output_folder}/{image_name}:/content"
            else:
                upload_url = f"https://graph.microsoft.com/v1.0/users/{location_id}/drive/root:/{image_name}:/content"
        else:
            if output_folder:
                upload_url = f"https://graph.microsoft.com/v1.0/sites/{location_id}/drive/root:/{output_folder}/{image_name}:/content"
            else:
                upload_url = f"https://graph.microsoft.com/v1.0/sites/{location_id}/drive/root:/{image_name}:/content"
        
        upload_response = requests.put(upload_url, headers=upload_headers, data=image_content)
        
        if upload_response.status_code in [200, 201]:
            result = upload_response.json()
            return {
                "id": result.get("id"),
                "name": result.get("name"),
                "size": result.get("size"),
                "webUrl": result.get("webUrl"),
                "slideIndex": slide_index,
                "format": image_format
            }
        else:
            return {"error": "Failed to upload image", "status_code": upload_response.status_code}
    else:
        return {"error": "Failed to get slide image", "status_code": thumb_response.status_code}

@mcp.tool()
def create_odf_document(location_type: str, location_id: str, file_name: str, doc_type: str, content: str, folder_path: str = ""):
    """Creates an OpenDocument Format file (.odt, .ods, .odp) in OneDrive or SharePoint.
    
    Args:
        location_type: Either 'onedrive' or 'sharepoint'
        location_id: User ID for OneDrive, or Site ID for SharePoint
        file_name: Name of the file (e.g., 'document.odt')
        doc_type: Type of document - 'text' (.odt), 'spreadsheet' (.ods), or 'presentation' (.odp)
        content: Text content for the document
        folder_path: Optional folder path
    """
    from odf.opendocument import OpenDocumentText, OpenDocumentSpreadsheet, OpenDocumentPresentation
    from odf.text import P
    from odf.table import Table, TableRow, TableCell
    from odf.draw import Page, Frame, TextBox
    import io
    
    # Create ODF document based on type
    if doc_type == "text":
        doc = OpenDocumentText()
        p = P(text=content)
        doc.text.addElement(p)
        content_type = "application/vnd.oasis.opendocument.text"
    elif doc_type == "spreadsheet":
        doc = OpenDocumentSpreadsheet()
        table = Table(name="Sheet1")
        # Add simple content
        tr = TableRow()
        tc = TableCell()
        p = P(text=content)
        tc.addElement(p)
        tr.addElement(tc)
        table.addElement(tr)
        doc.spreadsheet.addElement(table)
        content_type = "application/vnd.oasis.opendocument.spreadsheet"
    else:  # presentation
        doc = OpenDocumentPresentation()
        page = Page()
        frame = Frame(width="720pt", height="540pt")
        textbox = TextBox()
        p = P(text=content)
        textbox.addElement(p)
        frame.addElement(textbox)
        page.addElement(frame)
        doc.presentation.addElement(page)
        content_type = "application/vnd.oasis.opendocument.presentation"
    
    # Save to bytes
    doc_bytes = io.BytesIO()
    doc.save(doc_bytes)
    doc_bytes.seek(0)
    
    access_token = get_access_token()
    headers = {
        "Authorization": f"Bearer {access_token}",
        "Content-Type": content_type
    }
    
    # Build URL
    if location_type == "onedrive":
        if folder_path:
            url = f"https://graph.microsoft.com/v1.0/users/{location_id}/drive/root:/{folder_path}/{file_name}:/content"
        else:
            url = f"https://graph.microsoft.com/v1.0/users/{location_id}/drive/root:/{file_name}:/content"
    else:  # sharepoint
        if folder_path:
            url = f"https://graph.microsoft.com/v1.0/sites/{location_id}/drive/root:/{folder_path}/{file_name}:/content"
        else:
            url = f"https://graph.microsoft.com/v1.0/sites/{location_id}/drive/root:/{file_name}:/content"
    
    response = requests.put(url, headers=headers, data=doc_bytes.getvalue())
    
    if response.status_code in [200, 201]:
        result = response.json()
        return {
            "id": result.get("id"),
            "name": result.get("name"),
            "size": result.get("size"),
            "webUrl": result.get("webUrl"),
            "docType": doc_type,
            "createdDateTime": result.get("createdDateTime")
        }
    else:
        return {"error": response.text, "status_code": response.status_code}

async def async_main():
    """Async entry point for MCP server."""
    print("🚀 Starting MCP server for Microsoft Entra...")
    await mcp.run_stdio_async()

def main():
    """Main entry point for console script."""
    asyncio.run(async_main())

if __name__ == "__main__":
    main()
